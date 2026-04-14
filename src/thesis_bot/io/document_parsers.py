from __future__ import annotations

import io
import zipfile
from xml.etree import ElementTree

import fitz
from pptx import Presentation

from thesis_bot.io.document_source import DocumentArtifact, ParsedDocument


def parse_document_artifact(artifact: DocumentArtifact) -> ParsedDocument:
    """Parse a document artifact into normalized text."""
    extension = artifact.extension.lower()
    if extension == ".pdf":
        text = _extract_text_from_pdf_bytes(artifact.content)
    elif extension == ".docx":
        text = _extract_text_from_docx_bytes(artifact.content)
    elif extension == ".pptx":
        text = _extract_text_from_pptx_bytes(artifact.content)
    elif extension in {".md", ".txt"}:
        text = _decode_text_bytes(artifact.content)
    else:
        raise ValueError(f"Unsupported document extension: {artifact.extension}")

    return ParsedDocument(
        name=artifact.name,
        source_uri=artifact.source_uri,
        extension=artifact.extension,
        text=text,
    )


def parse_document_artifacts(artifacts: list[DocumentArtifact]) -> list[ParsedDocument]:
    """Parse a list of artifacts, skipping files that fail to parse."""
    parsed_documents: list[ParsedDocument] = []
    for artifact in artifacts:
        print(f"Extracting text from: {artifact.name}")
        try:
            document = parse_document_artifact(artifact)
            parsed_documents.append(document)
            print(f"  Loaded {len(document.text):,} characters from {artifact.extension}")
        except Exception as exc:
            print(f"  Failed to parse {artifact.name}: {exc}")
    return parsed_documents


def _extract_text_from_pdf_bytes(content: bytes) -> str:
    document = fitz.open(stream=content, filetype="pdf")
    try:
        text = ""
        for page_num in range(len(document)):
            text += document[page_num].get_text()
        return text
    finally:
        document.close()


def _extract_text_from_docx_bytes(content: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        xml_bytes = archive.read("word/document.xml")
    root = ElementTree.fromstring(xml_bytes)
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: list[str] = []
    for paragraph in root.findall(".//w:p", namespace):
        texts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
        joined = "".join(texts).strip()
        if joined:
            paragraphs.append(joined)
    return "\n\n".join(paragraphs)


def _extract_text_from_pptx_bytes(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    slides: list[str] = []
    for slide in presentation.slides:
        fragments: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                fragments.append(text.strip())
        if fragments:
            slides.append("\n".join(fragments))
    return "\n\n".join(slides)


def _decode_text_bytes(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("utf-8", errors="replace")
